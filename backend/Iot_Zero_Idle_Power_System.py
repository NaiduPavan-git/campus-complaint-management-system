from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Function to create a slide with title, bullets, and image placeholder
def create_content_slide(title, bullets, image_desc):
    # Use blank layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add Title
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(9)
    height = Inches(0.8)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(30)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 42, 84)  # Dark Blue

    # Add Bullet Points (Max 5)
    left_b = Inches(0.5)
    top_b = Inches(1.2)
    width_b = Inches(5.5)
    height_b = Inches(5.5)
    body_box = slide.shapes.add_textbox(left_b, top_b, width_b, height_b)
    tf = body_box.text_frame
    tf.word_wrap = True

    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_after = Pt(12)
    
    # Add Image Placeholder (User should replace this)
    # Ideally, you would use slide.shapes.add_picture('path/to/image.jpg', ...) here
    left_i = Inches(6.5)
    top_i = Inches(1.5)
    width_i = Inches(3.0)
    height_i = Inches(4.5)
    
    # Creating a visual placeholder box
    placeholder = slide.shapes.add_shape(1, left_i, top_i, width_i, height_i) # Rectangle
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(230, 230, 230)
    placeholder.line.color.rgb = RGBColor(100, 100, 100)
    
    # Text inside placeholder
    text_frame = placeholder.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = f"IMAGE PLACEHOLDER\n\n'{image_desc}'"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(80, 80, 80)

# --- SLIDE 1: TITLE ---
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "IoT Cloud-Zero Idle Power System"
subtitle.text = ("An IoT & Cloud-Based System to Detect and Eliminate Energy Wastage\n\n"
                 "Submitted By: Naidu Pavan (24VV1F0036)\n"
                 "Project Guide: Dr. Ch. Bindu Madhuri")

# --- SLIDE 2: PROBLEM STATEMENT ---
create_content_slide(
    "Problem Statement",
    [
        "Significant electricity wastage occurs when classrooms are empty but appliances remain ON.",
        "Primary causes include human negligence, cancelled classes, and lack of monitoring infrastructure.",
        "Existing automated solutions are often too expensive for large-scale institutional deployment.",
        "Manual supervision is inefficient and cannot provide real-time monitoring.",
        "Need for a low-cost, scalable system to quantify and alert on 'idle power' without forced cut-offs."
    ],
    "Empty classroom with lights left on"
)

# --- SLIDE 3: MOTIVATION & OBJECTIVES ---
create_content_slide(
    "Motivation & Objectives",
    [
        "Reduce idle power consumption using IoT devices and cloud-based analytics.",
        "Real-time detection of occupants and monitoring of appliance ON/OFF states.",
        "Calculate specific 'Idle Power Time' (room empty but devices ON) for behavioral analysis.",
        "Provide eco-feedback and historical data visualization via web dashboards.",
        "Maintain privacy (no cameras) and human-centric control (alerts vs. automatic cut-off)."
    ],
    "Concept of Eco-feedback or Green Energy"
)

# --- SLIDE 4: SYSTEM ARCHITECTURE ---
create_content_slide(
    "System Architecture",
    [
        "Perception Layer: ESP32 microcontroller with mmWave, Ultrasonic, and Current sensors.",
        "Network Layer: MQTT/HTTP protocols for real-time data transmission to the cloud.",
        "Cloud Layer: Rule-based engine to detect idle conditions and generate alerts.",
        "Database: Firebase for storing sensor logs and calculated energy metrics.",
        "Application Layer: Web dashboard for visualization and role-based access control."
    ],
    "IoT Architecture Diagram (3-Layer)"
)

# --- SLIDE 5: METHODOLOGY - HARDWARE ---
create_content_slide(
    "Methodology: Hardware Layer",
    [
        "Microcontroller: ESP32 for Wi-Fi connectivity and data processing.",
        "Occupancy Detection: Combination of mmWave and Ultrasonic sensors for robust accuracy.",
        "Energy Monitoring: ACS712 Current Sensor to estimate real-time power usage of fans/lights.",
        "Actuation: Relay modules for manual/automatic control of appliances.",
        "Local Indicators: LED/Buzzer alerts for immediate local feedback."
    ],
    "ESP32 Circuit Board with Sensors"
)

# --- SLIDE 6: METHODOLOGY - CLOUD & SOFTWARE ---
create_content_slide(
    "Methodology: Cloud & Software",
    [
        "Communication: MQTT protocol ensures low-latency data upload.",
        "Cloud Logic: Computes idle intervals and triggers alerts if wastage exceeds thresholds.",
        "Database: JSON-based storage of time-series sensor data (occupancy, device state).",
        "Dashboard Features: Live status maps, daily/weekly idle time graphs, and 'Energy Waste Leaderboards'.",
        "Notification Service: Automated Email/WhatsApp alerts sent to stakeholders."
    ],
    "Cloud Computing Server Icon / Code Snippet"
)

# --- SLIDE 7: RESULTS & ANALYTICS ---
create_content_slide(
    "Results & Behavioral Insights",
    [
        "Live Monitoring: Real-time identification of 'Empty but Powered ON' rooms (Red Alerts).",
        "Wastage Patterns: Data revealed peak energy waste during lunch hours (1 PM - 2 PM).",
        "Behavioral Impact: Introduction of the 'Wastage Leaderboard' reduced idle power by 15-20%.",
        "Room Ranking: Identification of high-traffic labs as frequent energy wasters.",
        "Alert Effectiveness: Successful delivery of notifications for idle power durations > 10 minutes."
    ],
    "Bar Chart showing Energy Wastage Statistics"
)

# --- SLIDE 8: CONCLUSION & FUTURE SCOPE ---
create_content_slide(
    "Conclusion & Future Scope",
    [
        "Cost-Effective: Scalable area-based sensing reduces hardware costs compared to per-node setups.",
        "Privacy Preserved: Non-invasive sensors avoid compromising student privacy.",
        "Future Scope 1: Integrate Machine Learning to predict occupancy patterns and optimize schedules.",
        "Future Scope 2: Sync with official college timetables to detect scheduled free periods automatically.",
        "Future Scope 3: Develop a dedicated mobile app for broader accessibility and alerts."
    ],
    "Smart City or Sustainable Future Concept"
)

# Save the Presentation
prs.save('IoT_Zero_Idle_Power_System.pptx')
print("Presentation 'IoT_Zero_Idle_Power_System.pptx' created successfully!")